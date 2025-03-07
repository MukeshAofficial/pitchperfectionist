
from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Dict, Any, List, Optional
import os
import tempfile
import uuid
import json
from pathlib import Path

# Try to import PowerPoint libraries - if they fail, provide guidance
try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Install with: pip install python-pptx")

try:
    import openai
except ImportError:
    print("WARNING: OpenAI not installed. Enhanced slide features will not work. Install with: pip install openai")

# Create FastAPI app
app = FastAPI(title="PitchPerfectionist API")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with your frontend origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Store presentations in memory (in production, use a database)
presentations = {}

class EnhanceSlideRequest(BaseModel):
    slide_index: int
    enhancement_type: str
    options: Dict[str, Any]

# Health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "ok"}

# Upload PowerPoint endpoint
@app.post("/upload-ppt")
async def upload_ppt(file: UploadFile = File(...)):
    # Check file extension
    if not file.filename.lower().endswith(('.ppt', '.pptx', '.doc', '.docx')):
        raise HTTPException(status_code=400, detail="File must be a PowerPoint or Word document")
    
    # Create a temporary file
    try:
        suffix = Path(file.filename).suffix
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            # Write the uploaded file content to the temporary file
            content = await file.read()
            temp_file.write(content)
            temp_path = temp_file.name
            
        # Process the file based on its type
        if suffix.lower() in ['.ppt', '.pptx']:
            # PowerPoint processing
            prs = Presentation(temp_path)
            slides = []
            
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                slides.append(slide_text.strip() or "No text on this slide")
        
        elif suffix.lower() in ['.doc', '.docx']:
            # Word document processing - simplified by treating each paragraph as a slide
            import docx
            doc = docx.Document(temp_path)
            slides = []
            current_slide = ""
            
            for para in doc.paragraphs:
                # If paragraph is a header, start a new slide
                if para.style.name.startswith('Heading'):
                    if current_slide:
                        slides.append(current_slide.strip())
                    current_slide = para.text + "\n"
                else:
                    current_slide += para.text + "\n"
                    
            # Add the last slide
            if current_slide:
                slides.append(current_slide.strip())
                
            # If no headers were found, treat the whole document as one slide
            if not slides:
                slides = ["".join(p.text + "\n" for p in doc.paragraphs)]
        
        # Generate a unique ID for the presentation
        presentation_id = str(uuid.uuid4())
        presentations[presentation_id] = {
            "filename": file.filename,
            "slides": slides
        }
        
        # Clean up the temporary file
        os.unlink(temp_path)
        
        return {
            "presentation_id": presentation_id,
            "filename": file.filename,
            "slides": slides,
            "slide_count": len(slides)
        }
    
    except Exception as e:
        # Clean up if there was an error
        if 'temp_path' in locals():
            os.unlink(temp_path)
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

# Get presentation slides
@app.get("/presentation/{presentation_id}")
async def get_presentation(presentation_id: str):
    if presentation_id not in presentations:
        raise HTTPException(status_code=404, detail="Presentation not found")
    
    return presentations[presentation_id]

# Enhance slide endpoint
@app.post("/enhance-slide")
async def enhance_slide(request: EnhanceSlideRequest, background_tasks: BackgroundTasks):
    try:
        # Get enhancement options
        slide_index = request.slide_index
        enhancement_type = request.enhancement_type
        options = request.options
        
        # For demonstration, we'll use the first presentation in memory
        if not presentations:
            raise HTTPException(status_code=404, detail="No presentations available")
        
        presentation_id = list(presentations.keys())[0]
        presentation = presentations[presentation_id]
        
        if slide_index >= len(presentation["slides"]):
            raise HTTPException(status_code=400, detail="Slide index out of range")
        
        slide_content = presentation["slides"][slide_index]
        
        # Check if we have OpenAI credentials
        api_key = options.get('apiKey')
        model = options.get('model', 'gpt-4o')
        
        if api_key:
            # Use OpenAI for enhancement
            enhanced_content = await enhance_with_openai(
                api_key, model, enhancement_type, slide_content, options
            )
        else:
            # No API key, use mock enhancement
            enhanced_content = mock_enhancement(enhancement_type, slide_content)
            
        return {
            "original_content": slide_content,
            "enhanced_content": enhanced_content,
            "enhancement_type": enhancement_type
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error enhancing slide: {str(e)}")

async def enhance_with_openai(api_key: str, model: str, enhancement_type: str, slide_content: str, options: Dict[str, Any]):
    try:
        openai.api_key = api_key
        
        # Create prompt based on enhancement type
        prompt = "Enhance the following presentation slide content"
        
        # Add enhancement type specific instructions
        enhancement_types = {
            "simplify": "by simplifying the language and making it easier to understand",
            "storytelling": "by adding storytelling elements and narrative structure",
            "professional": "by making it more professional and business-appropriate",
            "accessible": "by improving accessibility and inclusivity",
            "engaging": "by making it more engaging and captivating for the audience",
            "concise": "by making it more concise and focused on key points",
            "elaborate": "by elaborating on the content with more details and explanation",
            "academic": "by giving it an academic style appropriate for research or educational settings",
            "creative": "by adding creative elements and unique perspectives"
        }
        
        prompt += f" {enhancement_types.get(enhancement_type, '')}"
        
        # Add audience information if provided
        target_audience = options.get('targetAudience')
        if target_audience:
            prompt += f". The target audience is: {target_audience}"
        
        # Add tone level information
        tone_level = options.get('toneLevel', 5)
        prompt += f". Use a tone intensity of {tone_level}/10 (where 10 is the strongest)."
        
        # Add custom instructions if provided
        custom_instructions = options.get('customInstructions')
        if custom_instructions:
            prompt += f" Additional instructions: {custom_instructions}"
        
        # Add default prompt if set
        default_prompt = options.get('defaultPrompt')
        if default_prompt:
            prompt += f" {default_prompt}"
        
        # Add the slide content
        prompt += f"\n\nSlide content:\n{slide_content}"
        
        # Final instruction
        prompt += "\n\nEnhanced version:"
        
        # Call OpenAI API
        response = await openai.ChatCompletion.acreate(
            model=model,
            messages=[
                {"role": "system", "content": "You are an expert presentation writer and editor. Your task is to enhance presentation slide content based on specific requirements."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        return response.choices[0].message.content.strip()
    
    except Exception as e:
        print(f"OpenAI API error: {str(e)}")
        # Fallback to mock enhancement
        return mock_enhancement(enhancement_type, slide_content)

def mock_enhancement(enhancement_type: str, slide_content: str) -> str:
    """Generate mock enhanced content when API is not available"""
    enhancement_examples = {
        "simplify": f"{slide_content}\n\n[SIMPLIFIED VERSION]\nThis slide now uses clearer language and simpler explanations.",
        "storytelling": f"{slide_content}\n\n[STORYTELLING VERSION]\nThis slide now includes a compelling narrative arc and emotional elements.",
        "professional": f"{slide_content}\n\n[PROFESSIONAL VERSION]\nThis slide now uses business terminology and formal language.",
        "accessible": f"{slide_content}\n\n[ACCESSIBLE VERSION]\nThis slide now uses inclusive language and avoids jargon.",
        "engaging": f"{slide_content}\n\n[ENGAGING VERSION]\nThis slide now includes questions and interactive elements.",
        "concise": f"{slide_content}\n\n[CONCISE VERSION]\nThis slide now focuses on key messages without redundancy.",
        "elaborate": f"{slide_content}\n\n[ELABORATE VERSION]\nThis slide now includes more details and examples.",
        "academic": f"{slide_content}\n\n[ACADEMIC VERSION]\nThis slide now includes citations and research-based structure.",
        "creative": f"{slide_content}\n\n[CREATIVE VERSION]\nThis slide now incorporates metaphors and innovative presentation styles."
    }
    
    return enhancement_examples.get(enhancement_type, f"{slide_content}\n\n[ENHANCED VERSION]\nThis is a mock enhancement.")

# Run the server with: uvicorn server:app --reload
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="0.0.0.0", port=8000, reload=True)
